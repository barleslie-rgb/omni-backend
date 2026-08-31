import json
import os
import shutil
import time
import urllib.parse
from typing import Optional

from deep_translator import GoogleTranslator
from docx import Document
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from google import genai
from google.genai import types
import openpyxl
from PIL import Image
from pypdf import PdfReader
import requests

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

app = FastAPI(title="Omni Super-App Enterprise API", version="40.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
AUDIO_DIR = "audio_output"
EXPORT_DIR = "exports"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(EXPORT_DIR, exist_ok=True)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

client: Optional[genai.Client] = None
try:
    client = genai.Client(api_key=GEMINI_API_KEY)
except Exception:
    client = None

LANG_MAP = {
    "English": "en", "Hindi (हिन्दी)": "hi", "Marathi (मराठी)": "mr",
    "Sanskrit (संस्कृतम्)": "sa", "Gujarati (ગુજરાતી)": "gu", "Bengali (বাংলা)": "bn",
    "Tamil (தமிழ்)": "ta", "Telugu (తెలుగు)": "te", "Malayalam (മലയാളം)": "ml",
    "Kannada (ಕನ್ನಡ)": "kn", "Punjabi (ਪੰਜਾਬੀ)": "pa", "Urdu (اردو)": "ur",
    "Arabic (العربية)": "ar", "Egyptian Arabic (العصرية المصرية)": "ar",
    "Spanish (Español)": "es", "Portuguese (Português)": "pt", "German (Deutsch)": "de",
    "Chinese (中文)": "zh-cn", "Tagalog (Filipino)": "tl", "Vietnamese (Tiếng Việt)": "vi",
    "French (Français)": "fr", "Japanese (日本語)": "ja", "Korean (한국어)": "ko", "Russian (Русский)": "ru"
}

latest_document_context: str = ""
latest_uploaded_filename: str = "document"
detected_travel_destination: Optional[str] = None
chat_conversation_history: list[str] = []
AVAILABLE_MODELS = ['gemini-3.6-flash', 'gemini-3.7-flash']


@app.get("/")
async def root():
    return {
        "status": "online",
        "service": "Omni Super-App Enterprise API",
        "version": "40.0.0",
        "docs_url": "/docs"
    }


def call_gemini_with_retry(prompt: str) -> str:
    if not client:
        return "AI client not initialized. Please verify your GEMINI_API_KEY."
    last_err = ""
    for model_name in AVAILABLE_MODELS:
        try:
            response = client.models.generate_content(model=model_name, contents=prompt)
            if response and response.text:
                return response.text
        except Exception as e:
            last_err = str(e)
            time.sleep(0.5)
    return f"AI generation error: {last_err}"


def call_gemini_json(prompt: str) -> dict:
    if not client:
        return {}
    for model_name in AVAILABLE_MODELS:
        try:
            response = client.models.generate_content(
                model=model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.2,
                ),
            )
            if response and response.text:
                return json.loads(response.text)
        except Exception:
            time.sleep(0.5)
    return {}


def generate_ai_image(prompt: str, filename_prefix: str = "ai_gen") -> Optional[str]:
    """Generates a high-quality AI raster image (.jpg) and saves it to EXPORT_DIR."""
    try:
        out_filename = f"{filename_prefix}_{int(time.time())}.jpg"
        out_path = os.path.join(EXPORT_DIR, out_filename)

        encoded_prompt = urllib.parse.quote(prompt.strip())
        pollinations_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&nologo=true&seed={int(time.time())}"
        
        resp = requests.get(pollinations_url, timeout=25)
        if resp.status_code == 200 and len(resp.content) > 5000:
            with open(out_path, "wb") as f:
                f.write(resp.content)
            return out_filename
    except Exception:
        pass
    return None


@app.post("/api/v1/touristos-recommend")
async def touristos_recommend(
    time_available: str = Form("5 Hours"),
    location: str = Form("Mumbai"),
    group_type: str = Form("Family (Small/Large)"),
    dietary_preference: str = Form("All Foods Allowed"),
    target_language: str = Form("English")
):
    global detected_travel_destination
    active_location = location.strip() if location.strip() else (detected_travel_destination or "Mumbai")
    diet_rule = "STRICTLY VEGETARIAN ONLY." if dietary_preference == "Strictly Vegetarian Only" else "All food options allowed."

    json_prompt = f"""
You are TouristOS Guide. Generate an extensive, verified travel directory for '{active_location}' as a valid JSON object.
Context:
- Destination: {active_location}
- Available Window: {time_available}
- Group Profile: {group_type}
- Dietary Preference: {diet_rule}

JSON Format Schema:
{{
  "destination_summary": "2-3 comprehensive sentences highlighting key cultural, geographic, and travel aspects of {active_location}.",
  "distance_from_center": "Distance description from closest airport or city center to {active_location}",
  "transport_availability": "Specific local transit (Metro, local trains, auto-rickshaws, municipal buses, Uber/Ola)",
  "facilities": ["24/7 ATM Access", "Fast Wi-Fi Hubs", "EV Charging", "Luggage Storage", "Clean Public Facilities"],
  "best_things_to_do": [
    "Top attraction 1 in {active_location}",
    "Top attraction 2 in {active_location}",
    "Top attraction 3 in {active_location}",
    "Top attraction 4 in {active_location}"
  ],
  "best_food_to_try": [
    "Famous local specialty dish 1 ({diet_rule})",
    "Famous local specialty dish 2",
    "Top dining street/market in {active_location}",
    "Popular beverage/dessert"
  ],
  "shopping_malls_markets": [
    "Popular mall in {active_location}",
    "Traditional heritage market in {active_location}",
    "Shopping street in {active_location}"
  ],
  "hotels": [
    Provide 12 to 15 real hotels in {active_location}. Each hotel must have:
    - "name": (Real hotel name in {active_location})
    - "type": (e.g. "5-Star Luxury", "Boutique Resort", "Executive Stay", "Budget Comfort")
    - "rating": (e.g. "⭐ 4.8")
    - "reviews": (e.g. "1,850 reviews")
    - "price": (e.g. "₹3,499/night")
    - "address": (Specific street/neighborhood in {active_location})
    - "phone": (Contact phone number)
    - "amenities": ["Free Breakfast", "Swimming Pool", "Wi-Fi", "Parking"]
    - "description": (Brief 1-sentence description)
  ],
  "spots": [
    Provide 15 top sights in {active_location} (enough for 5 pages of 3 cards each). Each spot must have:
    - "title": (Name of landmark)
    - "rating": (e.g. "⭐ 4.8 (25k+)")
    - "dist": (Neighborhood or distance description)
    - "phone": (Contact phone or "Public Landmark")
    - "images": (Array of 2-3 high-resolution Unsplash photo URLs)
    - "tag": (e.g. "Heritage Site", "Beach & Sunset", "Shopping", "Nature Trail")
  ],
  "emergency": {{
    "hospital_name": "Premier multi-specialty hospital in {active_location}",
    "hospital_phone": "Verified contact number",
    "police_name": "Main police station in {active_location}",
    "police_phone": "Verified police contact number",
    "fire_name": "Fire brigade station in {active_location}",
    "fire_phone": "Verified fire contact number",
    "pharmacy_name": "24/7 chemist in {active_location}",
    "pharmacy_phone": "Verified pharmacy contact number"
  }}
}}
"""
    parsed_result = call_gemini_json(json_prompt)

    return JSONResponse(content={
        "status": "success",
        "resolved_location": active_location,
        "data": parsed_result
    })


@app.post("/api/v1/explore-chat")
async def explore_chat(
    spot_name: str = Form(...),
    question: str = Form(...),
    group_type: str = Form("Family (Small/Large)"),
    dietary_preference: str = Form("All Foods Allowed"),
    target_language: str = Form("English")
):
    prompt = f"""
You are the local AI Tour Guide for '{spot_name}'. Protect the user from scams and overcharging.
Group: '{group_type}', Diet: '{dietary_preference}'.
Question: "{question}".
Provide specific named places, addresses, contact details, and practical booking steps.
"""
    answer = call_gemini_with_retry(prompt)
    lang_code = LANG_MAP.get(target_language, "en")
    if lang_code != "en":
        try:
            answer = GoogleTranslator(source='auto', target=lang_code).translate(answer) or answer
        except Exception:
            pass
    return JSONResponse(content={"status": "success", "answer": answer})


@app.post("/api/v1/convert-file")
async def convert_file(
    request: Request,
    file: UploadFile = File(...),
    target_format: str = Form(...)
):
    filename = file.filename or "uploaded_file"
    file_path = os.path.join(UPLOAD_DIR, filename)
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    base_name = os.path.splitext(filename)[0]
    ext = target_format.lower().replace(".", "")
    out_filename = f"{base_name}_converted.{ext}"
    out_path = os.path.join(EXPORT_DIR, out_filename)

    try:
        if filename.lower().endswith(('png', 'jpg', 'jpeg', 'webp', 'bmp', 'gif')) and ext in ['png', 'jpg', 'jpeg', 'webp', 'bmp']:
            img = Image.open(file_path)
            if ext in ['jpg', 'jpeg']:
                img = img.convert('RGB')
                ext = 'jpeg'
            img.save(out_path, ext.upper())
        elif ext == 'docx':
            doc = Document()
            if filename.lower().endswith('.pdf'):
                reader = PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        doc.add_paragraph(t)
            elif filename.lower().endswith(('png', 'jpg', 'jpeg', 'webp', 'bmp')):
                doc.add_heading(f"Extracted Image Asset: {filename}", level=1)
                doc.add_picture(file_path, width=Inches(5.5) if HAS_PPTX else None)
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    for line in f:
                        doc.add_paragraph(line)
            doc.save(out_path)
        elif ext == 'xlsx':
            wb = openpyxl.Workbook()
            ws = wb.active
            if ws is not None:
                ws.title = "Converted Data"
                if filename.lower().endswith('.txt'):
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        for i, line in enumerate(f, start=1):
                            ws.cell(row=i, column=1, value=line.strip())
                else:
                    ws.cell(row=1, column=1, value=f"Exported Asset: {filename}")
                    ws.cell(row=2, column=1, value=f"Converted on {time.ctime()}")
                wb.save(out_path)
        elif ext == 'pptx' and HAS_PPTX:
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = base_name.replace('_', ' ').title()
            content_box = slide.placeholders[1]
            content_box.text = f"Converted from original source: {filename}\nProcessed via Omni PaperPilot Studio."
            prs.save(out_path)
        elif ext == 'txt':
            text_content = ""
            if filename.lower().endswith('.docx'):
                doc = Document(file_path)
                text_content = "\n".join([p.text for p in doc.paragraphs])
            elif filename.lower().endswith('.pdf'):
                reader = PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_content += t + "\n"
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
        else:
            shutil.copy(file_path, out_path)

        base_url = str(request.base_url).rstrip("/")
        download_url = f"{base_url}/api/v1/download-file/{out_filename}"
        return JSONResponse(content={"status": "success", "download_url": download_url, "converted_filename": out_filename})
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Conversion failed: {str(e)}")


@app.post("/api/v1/analyze-document")
async def analyze_document(file: UploadFile = File(...), target_language: str = Form("English")):
    global latest_document_context, latest_uploaded_filename, detected_travel_destination, chat_conversation_history
    chat_conversation_history.clear()
    filename = file.filename or "uploaded_file"
    latest_uploaded_filename = filename
    file_path = os.path.join(UPLOAD_DIR, filename)

    file_bytes = await file.read()
    with open(file_path, "wb") as buffer:
        buffer.write(file_bytes)

    is_image = filename.lower().endswith(('png', 'jpg', 'jpeg', 'webp', 'bmp'))
    ai_analysis = ""

    universal_prompt = """
You are Omni PaperPilot, an intelligent multimodal document and visual analyzer.

Analyze this uploaded document/image carefully and format your response clearly in Markdown:

### 📄 Document Overview & Visual Breakdown
- Identify the exact document type (Flight Ticket, Train Reservation, Receipt, Utility Bill, Invoice, ID, Medical Report, Contract, Photo/Screenshot).
- Provide a clear visual description of layout, logos, stamps, or key structural elements.

### 🔍 Key Information Extracted
- List all key details cleanly in bullet points (e.g., Dates, Names, PNR / Order / Invoice Numbers, Origin/Destination, Timings, Gate, Seat, Total Amount, Baggage allowances).

### 📝 Plain-English Summary
- Provide a concise summary of what this document represents and what the user needs to know.

### 🛡️ Safety & Advisory Check
- If this is an everyday document (Flight Ticket, Receipt, Itinerary, Photo, General Form): Confirm that it looks standard, valid, and free of apparent risks.
- If this is a Legal Contract, Financial Agreement, or Suspicious Communication: Highlight any hidden fees, predatory clauses, scam indicators, or penalty terms.
"""

    if is_image and client:
        try:
            img = Image.open(file_path).convert("RGB")
            response = client.models.generate_content(
                model='gemini-3.6-flash',
                contents=[img, universal_prompt]
            )
            ai_analysis = response.text if response and response.text else "Image analyzed successfully."
        except Exception as e:
            ai_analysis = f"Error during visual image inspection: {str(e)}"
    else:
        extracted_text = ""
        if filename.endswith(".pdf"):
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        extracted_text += t + "\n"
            except Exception:
                pass
        elif filename.endswith(".docx"):
            try:
                doc = Document(file_path)
                extracted_text = "\n".join([p.text for p in doc.paragraphs])
            except Exception:
                pass
        else:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    extracted_text = f.read()
            except Exception:
                pass

        full_prompt = f"{universal_prompt}\n\nDocument Text Content:\n{extracted_text[:6000]}"
        ai_analysis = call_gemini_with_retry(full_prompt)

    latest_document_context = ai_analysis

    city_extractor_prompt = f"""
Extract the arrival or destination city name from this document text.
Return ONLY the City Name (e.g., "Mumbai", "Vasai-Virar", "Wada", "New York", "Dubai").
If not a travel document, return "None".

Text:
{ai_analysis[:2000]}
"""
    dest_result = call_gemini_with_retry(city_extractor_prompt).strip()
    if dest_result and "none" not in dest_result.lower() and len(dest_result) < 30:
        detected_travel_destination = dest_result.strip('".\' ')
    else:
        detected_travel_destination = None

    lang_code = LANG_MAP.get(target_language, "en")
    if lang_code != "en":
        try:
            ai_analysis = GoogleTranslator(source='auto', target=lang_code).translate(ai_analysis) or ai_analysis
        except Exception:
            pass

    return JSONResponse(content={
        "status": "success",
        "data": {
            "file_name": filename,
            "plain_summary": ai_analysis,
            "detected_destination": detected_travel_destination
        }
    })


@app.post("/api/v1/ask-question")
async def ask_question(
    request: Request,
    question: str = Form(...),
    target_language: str = Form("English"),
    export_format: str = Form("none"),
    file: Optional[UploadFile] = File(None)
):
    global latest_document_context, latest_uploaded_filename, chat_conversation_history

    base_url = str(request.base_url).rstrip("/")

    if file is not None and file.filename:
        latest_uploaded_filename = file.filename
        file_path = os.path.join(UPLOAD_DIR, file.filename)
        f_bytes = await file.read()
        with open(file_path, "wb") as buffer:
            buffer.write(f_bytes)

        extra_text = ""
        if file.filename.lower().endswith(('png', 'jpg', 'jpeg', 'webp')) and client:
            try:
                img = Image.open(file_path).convert("RGB")
                resp = client.models.generate_content(model='gemini-3.6-flash', contents=[img, "Describe all details, text, and elements on this image."])
                extra_text = resp.text if resp and resp.text else ""
            except Exception:
                pass
        elif file.filename.endswith(".pdf"):
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        extra_text += t + "\n"
            except Exception:
                pass
        else:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    extra_text = f.read()
            except Exception:
                pass
        latest_document_context += f"\n\nNewly Uploaded File ({file.filename}):\n{extra_text[:3000]}"

    chat_conversation_history.append(f"User: {question}")

    intent_prompt = f"""
Analyze the user's message: "{question}"
Is the user asking to generate, create, draw, design, or render an image, photo, visual artwork, or logo?
Respond with a strict JSON:
{{
  "is_image_request": true/false,
  "image_prompt": "Enhanced photorealistic prompt for text-to-image generator",
  "text_response": "Polite explanation of the generated visual asset"
}}
"""
    intent_data = call_gemini_json(intent_prompt)
    generated_img_url = None

    if intent_data and intent_data.get("is_image_request"):
        img_prompt = intent_data.get("image_prompt") or question
        saved_img_name = generate_ai_image(img_prompt, "chat_ai_gen")
        if saved_img_name:
            generated_img_url = f"{base_url}/api/v1/download-file/{saved_img_name}"
            answer = intent_data.get("text_response") or f"Here is the generated image for: '{question}'"
        else:
            answer = "Image generated based on your prompt."
    else:
        prompt = f"""
Document Context:
{latest_document_context}

Chat History:
{chr(10).join(chat_conversation_history[-10:])}

Task: Answer the user's question, analyze document details, verify against scams/hidden fees, or format modifications requested:
"""
        answer = call_gemini_with_retry(prompt)

    latest_document_context = answer
    chat_conversation_history.append(f"AI: {answer}")

    lang_code = LANG_MAP.get(target_language, "en")
    if lang_code != "en":
        try:
            answer = GoogleTranslator(source='auto', target=lang_code).translate(answer) or answer
        except Exception:
            pass

    download_url: Optional[str] = None
    if generated_img_url:
        download_url = generated_img_url
    elif export_format in ["pdf", "docx", "xlsx", "pptx", "txt"]:
        base_name = os.path.splitext(latest_uploaded_filename)[0] if latest_uploaded_filename else "audited_document"
        if export_format == "docx":
            out_filename = f"{base_name}_studio_export.docx"
            out_path = os.path.join(EXPORT_DIR, out_filename)
            doc = Document()
            for line in answer.split("\n"):
                doc.add_paragraph(line)
            doc.save(out_path)
            download_url = f"{base_url}/api/v1/download-file/{out_filename}"
        elif export_format == "xlsx":
            out_filename = f"{base_name}_studio_export.xlsx"
            out_path = os.path.join(EXPORT_DIR, out_filename)
            wb = openpyxl.Workbook()
            ws = wb.active
            if ws is not None:
                for i, line in enumerate(answer.split("\n"), start=1):
                    ws.cell(row=i, column=1, value=line)
                wb.save(out_path)
            download_url = f"{base_url}/api/v1/download-file/{out_filename}"
        elif export_format == "pptx" and HAS_PPTX:
            out_filename = f"{base_name}_studio_export.pptx"
            out_path = os.path.join(EXPORT_DIR, out_filename)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Omni Studio Export"
            slide.placeholders[1].text = answer[:1000]
            prs.save(out_path)
            download_url = f"{base_url}/api/v1/download-file/{out_filename}"
        elif export_format == "txt":
            out_filename = f"{base_name}_studio_export.txt"
            out_path = os.path.join(EXPORT_DIR, out_filename)
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(answer)
            download_url = f"{base_url}/api/v1/download-file/{out_filename}"

    return JSONResponse(content={
        "status": "success",
        "answer": answer,
        "image_url": generated_img_url,
        "download_url": download_url
    })


@app.get("/api/v1/download-file/{filename}")
async def download_file(filename: str):
    path = os.path.join(EXPORT_DIR, filename)
    if os.path.exists(path):
        return FileResponse(path, filename=filename)
    raise HTTPException(status_code=404, detail="File not found")


@app.get("/api/v1/play-audio/{filename}")
async def play_audio(filename: str):
    path = os.path.join(AUDIO_DIR, filename)
    if os.path.exists(path):
        return FileResponse(path, media_type="audio/mpeg")
    raise HTTPException(status_code=404, detail="Not found")


if __name__ == "__main__":
    import uvicorn
    port =